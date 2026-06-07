"""Batch extract text from .docx/.pptx/.pdf files and save as raw markdown.

Usage: python extract_office_files.py <source_dir> <output_dir>

Each output file gets YAML frontmatter with:
  - source_url: (empty — fill in if from URL)
  - ingested: today's date
  - sha256: hex digest of body content (for drift detection on re-ingest)
  - original_file: the source filename

Requires: python-docx, python-pptx, pymupdf
Install: pip install python-docx python-pptx pymupdf
"""

import os, sys, hashlib
from datetime import date
from pathlib import Path


def sha256(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8")).hexdigest()[:16]


def extract_docx(fp: Path) -> str:
    from docx import Document
    doc = Document(str(fp))
    lines = []
    for p in doc.paragraphs:
        if p.text.strip():
            lines.append(p.text.strip())
    return "\n".join(lines)


def extract_pptx(fp: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(fp))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        lines.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        lines.append(row_text)
    return "\n".join(lines)


def extract_pdf(fp: Path) -> str:
    import fitz  # pymupdf
    doc = fitz.open(str(fp))
    lines = []
    for page in doc:
        text = page.get_text()
        if text.strip():
            lines.append(text.strip())
    doc.close()
    return "\n---\n".join(lines)


EXTRACTORS = {
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".pdf": extract_pdf,
}


def main():
    if len(sys.argv) != 3:
        print(f"Usage: python {sys.argv[0]} <source_dir> <output_dir>")
        sys.exit(1)

    src_dir = Path(sys.argv[1])
    out_dir = Path(sys.argv[2])
    out_dir.mkdir(parents=True, exist_ok=True)

    today = date.today().isoformat()
    results = []
    total_chars = 0

    files = sorted(
        f for f in src_dir.iterdir()
        if f.suffix.lower() in EXTRACTORS and f.is_file()
    )

    for fp in files:
        ext = fp.suffix.lower()
        try:
            text = EXTRACTORS[ext](fp)
            name = fp.stem
            safe_name = f"{name}.md"
            out_path = out_dir / safe_name
            h = sha256(text)

            content = (
                f"---\n"
                f"来源网址:\n"
                f"摄入日期: {today}\n"
                f"哈希值: {h}\n"
                f"原始文件: {fp.name}\n"
                f"---\n\n"
                f"# {name}\n\n"
                f"{text}\n"
            )
            out_path.write_text(content, encoding="utf-8")
            results.append((fp.name, len(text), safe_name, ""))
            total_chars += len(text)
        except Exception as e:
            results.append((fp.name, 0, "", str(e)))

    ok = sum(1 for r in results if not r[3])
    print(f"Processed: {ok}/{len(files)} files")
    for name, size, out, err in results:
        if err:
            print(f"  FAIL {name}: {err}")
        else:
            print(f"  OK   {name} ({size:,d} chars) → {out}")
    print(f"Total: {total_chars:,d} chars")


if __name__ == "__main__":
    main()
